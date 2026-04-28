import os
import re
import shutil
import subprocess
from datetime import datetime, timezone
from pathlib import Path

KORDOC_BIN = os.getenv("KORDOC_BIN", "kordoc")
KORDOC_EXTS = {".pdf", ".docx", ".xlsx", ".xls", ".hwp", ".hwpx"}

# Best-effort page-marker normalization. kordoc's output isn't fully specified
# across versions, so we detect a few common variants and emit a consistent
# `<!-- page: N -->` HTML comment. The chunker / downstream readers can grep
# for these to attach page numbers to chunks; for now we just count occurrences
# and surface page_count in the frontmatter.
_PAGE_MARKER_PATTERNS = [
    re.compile(r"^\s*<!--\s*page[:\s]+(\d+)\s*-->\s*$", re.IGNORECASE | re.MULTILINE),
    re.compile(r"^\s*-{3,}\s*page\s+(\d+)\s*-{3,}\s*$", re.IGNORECASE | re.MULTILINE),
    re.compile(r"^\s*Page\s+(\d+)\s*$", re.MULTILINE),
    re.compile(r"^\s*\[Page\s+(\d+)\]\s*$", re.IGNORECASE | re.MULTILINE),
    re.compile(r"^\s*\f\s*$", re.MULTILINE),  # form-feed (\x0c) used by some PDF→text converters
]


def _normalize_page_markers(body: str) -> tuple[str, int]:
    """Replace assorted page-marker syntaxes with `<!-- page: N -->`.

    Returns (normalized_body, page_count). page_count is 0 when no markers
    are detected — that's fine, we just omit page_count from the frontmatter.
    """
    page_idx = 0

    def _renumber(match: re.Match) -> str:
        nonlocal page_idx
        page_idx += 1
        # Prefer captured number when present (groups vary by pattern).
        n = match.group(1) if match.groups() else str(page_idx)
        return f"<!-- page: {n} -->"

    out = body
    for pat in _PAGE_MARKER_PATTERNS:
        out = pat.sub(_renumber, out)
    return out, page_idx


def _frontmatter(
    src: Path,
    rel_path: str,
    file_hash: str,
    fmt: str,
    page_count: int = 0,
) -> str:
    """L1 structural metadata. Consumed by m3 and stored in chunks.metadata."""
    folder = "/".join(rel_path.split("/")[:-1])
    filename = src.name
    try:
        st = src.stat()
        size = st.st_size
        mtime = datetime.fromtimestamp(st.st_mtime, tz=timezone.utc).isoformat()
    except OSError:
        size = -1
        mtime = ""
    lines = [
        "---",
        f"source_path: {rel_path}",
        f"source_hash: {file_hash}",
        f"format: {fmt}",
        f"folder_path: {folder}",
        f"filename: {filename}",
        f"size_bytes: {size}",
        f"mtime: {mtime}",
        f"converted_at: {datetime.now(timezone.utc).isoformat()}",
    ]
    if page_count > 0:
        lines.append(f"page_count: {page_count}")
    lines.extend(["---", "", ""])
    return "\n".join(lines)


def _run_kordoc(src: Path) -> str:
    bin_ = shutil.which(KORDOC_BIN) or KORDOC_BIN
    result = subprocess.run(
        [bin_, "--silent", str(src)],
        capture_output=True,
        text=True,
        timeout=300,
    )
    if result.returncode != 0:
        raise RuntimeError(f"kordoc failed: {result.stderr.strip()}")
    return result.stdout


def _run_pptx(src: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(src))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(r.text for r in para.runs).strip()
                    if text:
                        parts.append(text)
        parts.append("")
    return "\n".join(parts)


def convert(src: Path, rel_path: str, file_hash: str, out_dir: Path) -> tuple[Path, str]:
    """Convert document to Markdown. Returns (md_path, status)."""
    ext = src.suffix.lower()
    fmt = ext.lstrip(".")
    md_path = out_dir / (rel_path.rsplit(".", 1)[0] + ".md")
    md_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        if ext in KORDOC_EXTS:
            body = _run_kordoc(src)
            status = "ok"
        elif ext == ".pptx":
            body = _run_pptx(src)
            status = "ok-pptx-fallback"
        else:
            raise ValueError(f"unsupported ext: {ext}")
    except Exception as e:
        body = f"<!-- conversion failed: {e} -->\n"
        status = f"fail:{type(e).__name__}"

    body, page_count = _normalize_page_markers(body)
    md_path.write_text(
        _frontmatter(src, rel_path, file_hash, fmt, page_count) + body,
        encoding="utf-8",
    )
    return md_path, status
