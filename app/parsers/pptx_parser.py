import logging
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.parsers.base import BaseParser, ParseResult, ParsedPage

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    SUPPORTED_EXTENSIONS = [".pptx"]

    def parse(self, file_path: str) -> ParseResult:
        prs = Presentation(file_path)
        pages = []
        all_texts = []

        for slide_idx, slide in enumerate(prs.slides):
            parts = []

            for shape in slide.shapes:
                # 텍스트 프레임
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)

                # 표
                if shape.has_table:
                    table_md = self._table_to_markdown(shape.table)
                    if table_md:
                        parts.append(table_md)

                # 이미지
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_text = self._extract_image_ocr(shape, slide_idx + 1)
                    if img_text:
                        parts.append(img_text)

                # 그룹 shape 내부 탐색
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in shape.shapes:
                        if hasattr(child, "text") and child.text.strip():
                            parts.append(child.text.strip())
                        if hasattr(child, "has_table") and child.has_table:
                            table_md = self._table_to_markdown(child.table)
                            if table_md:
                                parts.append(table_md)

            if parts:
                slide_text = f"## Slide {slide_idx + 1}\n\n" + "\n\n".join(parts)
                pages.append(ParsedPage(page_num=slide_idx + 1, text=slide_text))
                all_texts.append(slide_text)

        full_text = "\n\n".join(all_texts)

        return ParseResult(
            pages=pages,
            raw_text=full_text,
            total_pages=len(pages),
            metadata={
                "parser": "python-pptx",
                "source": file_path,
                "slide_count": len(prs.slides),
            },
        )

    def _table_to_markdown(self, table) -> str:
        """pptx Table → 마크다운 표."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        col_count = len(rows[0])
        lines = []
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("| " + " | ".join(["---"] * col_count) + " |")
        for row in rows[1:]:
            while len(row) < col_count:
                row.append("")
            lines.append("| " + " | ".join(row[:col_count]) + " |")

        return "\n".join(lines)

    def _extract_image_ocr(self, shape, slide_num: int) -> str:
        """슬라이드 이미지 → MinerU OCR."""
        try:
            image = shape.image
            blob = image.blob

            with tempfile.NamedTemporaryFile(
                suffix=f".{image.content_type.split('/')[-1]}",
                delete=False,
            ) as tmp:
                tmp.write(blob)
                tmp_path = tmp.name

            from app.parsers.image_parser import ImageParser

            parser = ImageParser()
            result = parser.parse(tmp_path)
            Path(tmp_path).unlink(missing_ok=True)

            if result.raw_text.strip():
                return result.raw_text.strip()
        except Exception as e:
            logger.debug("OCR failed for image on slide %d: %s", slide_num, e)

        return f"[슬라이드 {slide_num} 이미지]"
