"""Document parser module.

Delegates to MinerU API for PDF/image parsing, falls back to local parsers
for DOCX, XLSX, PPTX.
"""

import logging
import re
import tempfile
from dataclasses import dataclass, field
from functools import lru_cache
from pathlib import Path

import httpx

from rag_pipeline.config import pipeline_settings

logger = logging.getLogger(__name__)

CAPTION_RE = re.compile(r"^\s*(figure|fig\.|table|표|그림)\s*[\dA-Za-z\-.:)]*", re.IGNORECASE)


@dataclass
class ExtractedImage:
    temp_path: str
    page_num: int | None
    image_type: str
    width: int | None = None
    height: int | None = None


@dataclass
class ParseBlock:
    block_type: str
    source_text: str
    page_number: int | None = None
    sheet_name: str | None = None
    slide_number: int | None = None
    section_path: str | None = None
    parent_local_idx: int | None = None
    language: str = "ko"
    bbox: dict | None = None
    metadata: dict = field(default_factory=dict)


@dataclass
class ParseResult:
    raw_text: str
    blocks: list[ParseBlock] = field(default_factory=list)
    images: list[ExtractedImage] = field(default_factory=list)
    total_pages: int = 0
    metadata: dict = field(default_factory=dict)


def _compose_raw_text(blocks: list[ParseBlock]) -> str:
    return "\n\n".join(
        block.source_text.strip()
        for block in blocks
        if block.source_text.strip() and block.block_type in {"text", "table", "caption"}
    ).strip()


def _make_block(block_type: str, source_text: str, **kwargs) -> ParseBlock | None:
    text = (source_text or "").strip()
    if not text:
        return None
    return ParseBlock(block_type=block_type, source_text=text, **kwargs)


def _append_image_blocks(
    blocks: list[ParseBlock],
    images: list[ExtractedImage],
    *,
    label_prefix: str = "Extracted image",
    page_fallback: int | None = None,
    slide_number: int | None = None,
) -> None:
    for idx, image in enumerate(images, start=1):
        page_number = image.page_num or page_fallback
        block = _make_block(
            "image",
            f"{label_prefix} {idx} from page {page_number or 1}",
            page_number=page_number,
            slide_number=slide_number,
            metadata={"image_index": idx, "image_type": image.image_type},
        )
        if block:
            blocks.append(block)


def _link_image_blocks_to_captions(blocks: list[ParseBlock]) -> None:
    caption_positions = [
        (idx, block)
        for idx, block in enumerate(blocks)
        if block.block_type == "caption"
    ]
    if not caption_positions:
        return

    for idx, block in enumerate(blocks):
        if block.block_type != "image":
            continue

        candidates = [
            (caption_idx, caption)
            for caption_idx, caption in caption_positions
            if (
                caption.page_number == block.page_number
                or (caption.slide_number and caption.slide_number == block.slide_number)
            )
        ]
        if not candidates:
            continue

        caption_idx, caption = min(candidates, key=lambda item: abs(item[0] - idx))
        block.parent_local_idx = caption_idx
        block.metadata = {
            **(block.metadata or {}),
            "caption_text": caption.source_text,
            "caption_block_idx": caption_idx,
        }
        if caption.source_text and caption.source_text not in block.source_text:
            block.source_text = f"{block.source_text}\nCaption: {caption.source_text}".strip()


def _parse_markdown_blocks(markdown: str, page_number: int | None = None) -> list[ParseBlock]:
    blocks: list[ParseBlock] = []
    lines = markdown.splitlines()
    buffer: list[str] = []
    buffer_type = "text"

    def flush() -> None:
        nonlocal buffer, buffer_type
        if not buffer:
            return
        text = "\n".join(buffer).strip()
        block = _make_block(buffer_type, text, page_number=page_number)
        if block:
            blocks.append(block)
        buffer = []
        buffer_type = "text"

    for line in lines:
        stripped = line.strip()
        if not stripped:
            flush()
            continue

        if stripped.startswith("|") and stripped.endswith("|"):
            if buffer_type != "table":
                flush()
                buffer_type = "table"
            buffer.append(stripped)
            continue

        if CAPTION_RE.match(stripped):
            flush()
            block = _make_block("caption", stripped, page_number=page_number)
            if block:
                blocks.append(block)
            continue

        if stripped.startswith("#"):
            flush()
            block = _make_block("text", stripped, page_number=page_number, section_path=stripped.lstrip("# ").strip())
            if block:
                blocks.append(block)
            continue

        if buffer_type != "text":
            flush()
        buffer.append(stripped)

    flush()
    return blocks


# ---------------------------------------------------------------------------
# MinerU API client (PDF + image OCR)
# ---------------------------------------------------------------------------

@lru_cache(maxsize=1)
def _get_http_client() -> httpx.Client:
    return httpx.Client()


MINERU_EXTENSIONS = {".pdf", ".png", ".jpg", ".jpeg", ".tif", ".tiff"}
LOCAL_EXTENSIONS = {".docx", ".xlsx", ".xls", ".pptx"}


def parse_document(file_path: str) -> ParseResult:
    """Parse a document, routing to the appropriate parser."""
    ext = Path(file_path).suffix.lower()

    if ext in MINERU_EXTENSIONS:
        return _parse_via_mineru(file_path, ext)
    elif ext == ".docx":
        return _parse_docx(file_path)
    elif ext in (".xlsx", ".xls"):
        return _parse_xlsx(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ---------------------------------------------------------------------------
# MinerU API
# ---------------------------------------------------------------------------

def _parse_via_mineru(file_path: str, ext: str) -> ParseResult:
    """Call the MinerU API microservice."""
    method = "ocr" if ext in (".png", ".jpg", ".jpeg", ".tif", ".tiff") else "auto"
    timeout = 300.0 if method == "ocr" else 600.0

    response = _get_http_client().post(
        f"{pipeline_settings.mineru_api_url}/parse",
        json={
            "file_path": file_path,
            "method": method,
            "backend": pipeline_settings.mineru_backend,
            "lang": pipeline_settings.mineru_lang,
        },
        timeout=timeout,
    )

    if response.status_code != 200:
        detail = response.json().get("detail", response.text)
        raise RuntimeError(f"MinerU API error for {file_path}: {detail}")

    data = response.json()

    images = [
        ExtractedImage(
            temp_path=img["path"],
            page_num=img.get("page_num"),
            image_type=Path(img["path"]).suffix.lstrip("."),
        )
        for img in data.get("images", [])
    ]

    blocks: list[ParseBlock] = []
    pages = data.get("pages") or []
    if pages:
        for page in pages:
            page_no = page.get("page_num")
            page_text = page.get("text", "")
            blocks.extend(_parse_markdown_blocks(page_text, page_number=page_no))
    else:
        blocks.extend(_parse_markdown_blocks(data["markdown"]))

    _append_image_blocks(blocks, images)

    _link_image_blocks_to_captions(blocks)

    return ParseResult(
        raw_text=_compose_raw_text(blocks) or data["markdown"],
        blocks=blocks,
        images=images,
        total_pages=data["total_pages"],
        metadata=data.get("metadata", {}),
    )


# ---------------------------------------------------------------------------
# Local parsers (DOCX, XLSX, PPTX)
# ---------------------------------------------------------------------------

def _parse_docx(file_path: str) -> ParseResult:
    """Parse DOCX using python-docx."""
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn

    doc = DocxDocument(file_path)
    extracted_images: list[ExtractedImage] = []
    parts = []
    blocks: list[ParseBlock] = []
    image_count = 0
    seen_docx_image_rel_ids: set[str] = set()

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            runs_text = []
            for run in element.iter(qn("w:r")):
                t = run.find(qn("w:t"))
                if t is not None and t.text:
                    runs_text.append(t.text)
                drawing = run.find(qn("w:drawing"))
                if drawing is not None:
                    image_count += 1
                    # Extract image blob if possible
                    img = _extract_docx_inline_image(
                        doc,
                        drawing,
                        image_count,
                        extracted_images,
                        seen_docx_image_rel_ids,
                    )
                    if img:
                        runs_text.append(img)

            text = "".join(runs_text).strip()
            if text:
                parts.append(text)
                block_type = "caption" if CAPTION_RE.match(text) else "text"
                block = _make_block(block_type, text, page_number=1)
                if block:
                    blocks.append(block)

        elif tag == "tbl":
            table_md = _docx_table_to_markdown(element, qn)
            if table_md:
                parts.append(table_md)
                block = _make_block("table", table_md, page_number=1)
                if block:
                    blocks.append(block)

    # Extract all embedded images from relationships
    for rel_id, rel in doc.part.rels.items():
        if "image" in rel.reltype:
            if rel_id in seen_docx_image_rel_ids:
                continue
            try:
                blob = rel.target_part.blob
                ext = "png"
                with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                    tmp.write(blob)
                    extracted_images.append(ExtractedImage(
                        temp_path=tmp.name,
                        page_num=1,
                        image_type=ext,
                    ))
                    seen_docx_image_rel_ids.add(rel_id)
            except Exception as e:
                logger.debug("Failed to extract image from relationship: %s", e)

    _append_image_blocks(blocks, extracted_images, label_prefix="DOCX image", page_fallback=1)

    full_text = "\n\n".join(parts)

    if not full_text.strip():
        full_text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    _link_image_blocks_to_captions(blocks)

    return ParseResult(
        raw_text=_compose_raw_text(blocks) or full_text,
        blocks=blocks,
        images=extracted_images,
        total_pages=1,
        metadata={"parser": "python-docx", "source": file_path, "image_count": image_count},
    )


def _extract_docx_inline_image(doc, drawing_element, img_idx, extracted_images, seen_rel_ids):
    """Try to extract inline image from a drawing element."""
    from docx.oxml.ns import qn

    try:
        blip = drawing_element.find(".//" + qn("a:blip"))
        if blip is None:
            return f"[image {img_idx}]"
        embed_id = blip.get(qn("r:embed"))
        if not embed_id:
            return f"[image {img_idx}]"
        rel = doc.part.rels.get(embed_id)
        if rel is None:
            return f"[image {img_idx}]"
        if embed_id in seen_rel_ids:
            return f"[image {img_idx}]"
        blob = rel.target_part.blob
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(blob)
            extracted_images.append(ExtractedImage(
                temp_path=tmp.name,
                page_num=1,
                image_type="png",
            ))
        seen_rel_ids.add(embed_id)
        return f"[image {img_idx}]"
    except Exception as e:
        logger.debug("Failed to extract inline image %d: %s", img_idx, e)
        return f"[image {img_idx}]"


def _docx_table_to_markdown(tbl_element, qn) -> str:
    """Convert lxml table element to markdown table."""
    rows = []
    for tr in tbl_element.iter(qn("w:tr")):
        cells = []
        for tc in tr.iter(qn("w:tc")):
            cell_texts = []
            for p in tc.iter(qn("w:p")):
                t_parts = []
                for t in p.iter(qn("w:t")):
                    if t.text:
                        t_parts.append(t.text)
                cell_texts.append("".join(t_parts))
            cells.append(" ".join(cell_texts).strip())
        rows.append(cells)

    if not rows:
        return ""

    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    for row in rows[1:]:
        while len(row) < len(rows[0]):
            row.append("")
        lines.append("| " + " | ".join(row[:len(rows[0])]) + " |")

    return "\n".join(lines)


def _parse_xlsx(file_path: str) -> ParseResult:
    """Parse XLSX/XLS using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    all_texts = []
    blocks: list[ParseBlock] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(c) if c is not None else "" for c in row]
            if any(v.strip() for v in cell_values):
                rows.append(" | ".join(cell_values))

        if rows:
            sheet_text = f"## Sheet: {sheet_name}\n\n" + "\n".join(rows)
            all_texts.append(sheet_text)
            block = _make_block("table", sheet_text, page_number=1, sheet_name=sheet_name, section_path=f"Sheet:{sheet_name}")
            if block:
                blocks.append(block)

    wb.close()
    full_text = "\n\n".join(all_texts)

    return ParseResult(
        raw_text=_compose_raw_text(blocks) or full_text,
        blocks=blocks,
        total_pages=len(all_texts) or 1,
        metadata={"parser": "openpyxl", "source": file_path},
    )


def _parse_pptx(file_path: str) -> ParseResult:
    """Parse PPTX using python-pptx."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    extracted_images: list[ExtractedImage] = []
    all_texts = []
    blocks: list[ParseBlock] = []
    image_counter = 0

    for slide_idx, slide in enumerate(prs.slides):
        parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
                    block_type = "caption" if CAPTION_RE.match(text) else "text"
                    block = _make_block(block_type, text, slide_number=slide_idx + 1, page_number=slide_idx + 1)
                    if block:
                        blocks.append(block)

            if shape.has_table:
                table_md = _pptx_table_to_markdown(shape.table)
                if table_md:
                    parts.append(table_md)
                    block = _make_block("table", table_md, slide_number=slide_idx + 1, page_number=slide_idx + 1)
                    if block:
                        blocks.append(block)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    blob = image.blob
                    ext = image.content_type.split("/")[-1]
                    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                        tmp.write(blob)
                        image_counter += 1
                        extracted_images.append(ExtractedImage(
                            temp_path=tmp.name,
                            page_num=slide_idx + 1,
                            image_type=ext,
                        ))
                        image_block = _make_block(
                            "image",
                            f"Slide {slide_idx + 1} image",
                            slide_number=slide_idx + 1,
                            page_number=slide_idx + 1,
                            metadata={"image_type": ext, "image_index": image_counter},
                        )
                        if image_block:
                            blocks.append(image_block)
                except Exception as e:
                    logger.debug("Failed to extract image on slide %d: %s", slide_idx + 1, e)

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if hasattr(child, "text") and child.text.strip():
                        parts.append(child.text.strip())
                        block_type = "caption" if CAPTION_RE.match(child.text.strip()) else "text"
                        block = _make_block(block_type, child.text.strip(), slide_number=slide_idx + 1, page_number=slide_idx + 1)
                        if block:
                            blocks.append(block)
                    if hasattr(child, "has_table") and child.has_table:
                        table_md = _pptx_table_to_markdown(child.table)
                        if table_md:
                            parts.append(table_md)
                            block = _make_block("table", table_md, slide_number=slide_idx + 1, page_number=slide_idx + 1)
                            if block:
                                blocks.append(block)

        if parts:
            slide_text = f"## Slide {slide_idx + 1}\n\n" + "\n\n".join(parts)
            all_texts.append(slide_text)

    full_text = "\n\n".join(all_texts)

    _link_image_blocks_to_captions(blocks)

    return ParseResult(
        raw_text=_compose_raw_text(blocks) or full_text,
        blocks=blocks,
        images=extracted_images,
        total_pages=len(prs.slides),
        metadata={"parser": "python-pptx", "source": file_path, "slide_count": len(prs.slides)},
    )


def _pptx_table_to_markdown(table) -> str:
    """Convert pptx Table to markdown."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    col_count = len(rows[0])
    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * col_count) + " |")
    for row in rows[1:]:
        while len(row) < col_count:
            row.append("")
        lines.append("| " + " | ".join(row[:col_count]) + " |")

    return "\n".join(lines)
